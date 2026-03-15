import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define common styles
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NS-IR Compiler: Learned Cost Modeling"
    subtitle.text = "A Summary of Neural-Symbolic Intermediate Representation\n\nProject Overview"

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Problem: Code Optimization"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Optimizing code for modern hardware is challenging:"
    p = tf.add_paragraph()
    p.text = "Traditional compilers (LLVM, GCC) rely on human-written heuristics that often fail to generalize to new architectures."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Modern Auto-tuners trial thousands of code versions directly on hardware, which is extremely slow."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Goal: Create a fast, accurate AI model that can predict how fast code will run before actually compiling and executing it."
    p.level = 0

    # --- Slide 3: The NS-IR Solution ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Solution: NS-IR"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Neural-Symbolic Intermediate Representation (NS-IR):"
    p = tf.add_paragraph()
    p.text = "Instead of relying on human rules, we train a machine learning model to understand programs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "We extract the structure of the code (loops, operations, flow) into a mathematical graph."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A Deep Learning Transformer model analyzes this graph to instantly predict execution speedups."
    p.level = 1

    # --- Slide 4: Key Framework Upgrades ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Technical Enhancements"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "The system was recently upgraded to a state-of-the-art architecture:"
    p = tf.add_paragraph()
    p.text = "Advanced Attention: The model learns precisely how different optimizations (like loop unrolling) affect specific parts of the code."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Confidence Estimation: The model predicts not just the speedup, but its own certainty level, allowing for risk-aware scheduling."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Robust Training Pipeline: Trained on 50,000 code samples using advanced loss functions to ensure high accuracy."
    p.level = 1

    # --- Slide 5: Final Results ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Evaluation & Final Results"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "The upgraded model demonstrates exceptional predictive accuracy:"
    p = tf.add_paragraph()
    p.text = "Accuracy: Achieved a Mean Absolute Percentage Error (MAPE) of 2.28% on validation data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Speed: Inference latency is under 5 milliseconds per prediction."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Scale: Successfully converged over 5.8 million parameters, providing a highly capable replacement for traditional compiler models."
    p.level = 1

    # --- Slide 6: Conclusion ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Summary:"
    p = tf.add_paragraph()
    p.text = "The NS-IR Compiler bridges the gap between deep learning and traditional software compilation."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "By accurately and rapidly predicting execution costs, it paves the way for fully automated, AI-driven code optimization."
    p.level = 1

    output_filename = 'NS-IR_Summary_Presentation.pptx'
    prs.save(output_filename)
    print(f"Presentation created successfully as '{output_filename}'")

if __name__ == '__main__':
    create_presentation()
